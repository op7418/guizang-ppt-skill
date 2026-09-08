#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
html2pptx — 把 guizang-ppt-skill 生成的单文件 HTML deck 转成 .pptx（截图组装版）

用法:
    python scripts/html2pptx.py <index.html> [-o 输出.pptx] [--workers 8]

依赖:
    pip install python-pptx
    系统自带 Microsoft Edge 或 Google Chrome（自动探测，无需安装浏览器）

产出:
    与 index.html 同目录的 .pptx（16:9，每页 1920x1080 高清截图）

断点续跑:
    截图缓存于 <index.html 同目录>/_html2pptx_render/，
    重复执行只补缺失/损坏的页，秒级完成。

已知坑（脚本已全部自动规避，勿手改）:
    1. 模板 HTML 注释里含假 <section class="slide ..."> 字样 → 提取前剥离全部注释
    2. 隔离显示当前页后，再对 #deck 做 translateX 平移会把唯一可见页推出视口
       （display:none 的页不占 flex 位） → 只用 display 隔离，禁止任何 transform
    3. --virtual-time-budget < 5000ms 会截到空白页 → 固定 6000ms
    4. 并行多实例共用默认 profile 会锁冲突 → 每实例独立临时 --user-data-dir
    5. 空白截图特征是 ~8KB 纯色 → 以 20000 字节为有效阈值，低于自动重试 3 次
"""

import argparse
import glob
import os
import re
import shutil
import subprocess
import sys
import tempfile
import urllib.parse
from concurrent.futures import ThreadPoolExecutor

# ---------------- 可调参数 ----------------
W, H = 1920, 1080          # 截图分辨率
BUDGET_MS = 6000           # 虚拟时间预算（≥5000，低了截空白）
MIN_BYTES = 20000          # 有效截图的体积阈值（空白图约 8KB）
DEFAULT_WORKERS = 8

CAPTURE_TPL = """<!DOCTYPE html>
<html lang="zh-CN">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=1920">
<title>html2pptx capture</title>
__STYLE__
<style>
  /* capture 专用覆盖：静止 + 单页 + 满幅 */
  html,body{width:__W__px!important;height:__H__px!important;overflow:hidden}
  .bg{display:none!important}
  #nav,#hint{display:none!important}
  [data-anim]{opacity:1!important;transform:none!important;animation:none!important;transition:none!important}
</style>
</head>
<body>
<div id="deck">__SLIDES__</div>
<script>
  try{document.body.classList.add('low-power')}catch(e){}
  var q=new URLSearchParams(location.search),
      p=Math.max(1,Math.min(__TOTAL__,parseInt(q.get('p')||'1',10))),
      dots=document.querySelectorAll('#deck .slide');
  /* 坑2：只用 display 隔离。绝不对 #deck 做 translateX——
     display:none 不占 flex 位，当前页必在首位，平移只会把它推出视口。 */
  dots.forEach(function(s,i){s.style.display=(i===p-1)?'':'none'});
  var cur=dots[p-1];
  if(cur){
    var isLight=cur.classList.contains('light');
    document.body.style.background=isLight?'#f1f3f5':'#0a1f3d';
  }
</script>
</body>
</html>
"""


def find_browser():
    """探测 Edge / Chrome，返回可执行文件路径。"""
    cands = [
        os.path.expandvars(r"%ProgramFiles(x86)%\Microsoft\Edge\Application\msedge.exe"),
        os.path.expandvars(r"%ProgramFiles%\Microsoft\Edge\Application\msedge.exe"),
        os.path.expandvars(r"%LocalAppData%\Google\Chrome\Application\chrome.exe"),
        os.path.expandvars(r"%ProgramFiles%\Google\Chrome\Application\chrome.exe"),
        os.path.expandvars(r"%ProgramFiles(x86)%\Google\Chrome\Application\chrome.exe"),
        "/usr/bin/google-chrome", "/usr/bin/chromium-browser", "/usr/bin/chromium",
        "/Applications/Google Chrome.app/Contents/MacOS/Google Chrome",
    ]
    for c in cands:
        if c and os.path.exists(c):
            return c
    sys.exit("[html2pptx] 未找到 Edge/Chrome，请先安装浏览器")


def read_text(p):
    with open(p, "r", encoding="utf-8") as f:
        return f.read()


def build_capture(index_path, capture_path):
    """从 index.html 生成单页定位的 capture.html。"""
    html = read_text(index_path)
    # 坑1：剥离全部 HTML 注释（模板注释里有假 section 标签，会污染提取并吞掉真实首页）
    html = re.sub(r"<!--.*?-->", "", html, flags=re.S)

    styles = re.findall(r"<style[^>]*>.*?</style>", html, re.S)
    if not styles:
        sys.exit("[html2pptx] index.html 里找不到 <style> 块")

    slides = re.findall(r'<section class="slide.*?</section>', html, re.S)
    if not slides:
        sys.exit("[html2pptx] index.html 里找不到 slide section")

    # 健全性校验：每页必须带 light/dark 主题类（拦截注释污染 / 垃圾块）
    for i, s in enumerate(slides):
        m = re.match(r'<section class="slide([^"]*)"', s)
        cls = m.group(1).split() if m else []
        if not ({"light", "dark"} & set(cls)):
            sys.exit("[html2pptx] 第 %d 页 class 异常（缺 light/dark）：%s"
                     % (i + 1, m.group(1) if m else "?"))

    tpl = (CAPTURE_TPL
           .replace("__STYLE__", "\n".join(styles))
           .replace("__SLIDES__", "\n".join(slides))
           .replace("__TOTAL__", str(len(slides)))
           .replace("__W__", str(W))
           .replace("__H__", str(H)))
    with open(capture_path, "w", encoding="utf-8") as f:
        f.write(tpl)
    return len(slides)


def render_all(browser, capture_path, render_dir, total, workers):
    os.makedirs(render_dir, exist_ok=True)
    # 清理历史坏图（体积过小的空白截图）
    for f in glob.glob(os.path.join(render_dir, "*.png")):
        try:
            if os.path.getsize(f) < MIN_BYTES:
                os.remove(f)
        except OSError:
            pass

    url = "file:///" + urllib.parse.quote(os.path.abspath(capture_path).replace("\\", "/"))

    def one_shot(p):
        out = os.path.join(render_dir, "p%03d.png" % p)
        if os.path.exists(out) and os.path.getsize(out) > MIN_BYTES:
            return p, True  # 好图直接跳过（断点续跑）
        prof = os.path.join(tempfile.gettempdir(), "h2p_%d_%d" % (os.getpid(), p))
        cmd = [browser, "--headless", "--disable-gpu", "--hide-scrollbars",
               "--user-data-dir=" + prof,           # 坑4：每实例独立 profile
               "--enable-unsafe-swiftshader",
               "--window-size=%d,%d" % (W, H),
               "--virtual-time-budget=%d" % BUDGET_MS,  # 坑3：固定预算
               "--screenshot=" + out,
               "%s?p=%d" % (url, p)]
        ok = False
        for _ in range(3):                            # 坑5：坏图自动重试
            try:
                subprocess.run(cmd, capture_output=True, timeout=90)
            except subprocess.TimeoutExpired:
                pass
            if os.path.exists(out) and os.path.getsize(out) > MIN_BYTES:
                ok = True
                break
        shutil.rmtree(prof, ignore_errors=True)
        return p, ok

    with ThreadPoolExecutor(max_workers=workers) as ex:
        results = list(ex.map(one_shot, range(1, total + 1)))
    fails = [p for p, ok in results if not ok]
    if fails:
        sys.exit("[html2pptx] 以下页截图失败（重试3次仍无效）：%s" % fails)


def build_pptx(render_dir, out_path, total):
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError:
        sys.exit("[html2pptx] 缺少 python-pptx，请执行: pip install python-pptx")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    added = 0
    for p in range(1, total + 1):
        img = os.path.join(render_dir, "p%03d.png" % p)
        if not os.path.exists(img):
            continue
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(img, 0, 0, Inches(13.333), Inches(7.5))
        added += 1
    if added != total:
        sys.exit("[html2pptx] PPTX 页数 %d != 截图数 %d" % (added, total))
    prs.save(out_path)
    return added


def main():
    ap = argparse.ArgumentParser(description="HTML deck -> PPTX (screenshot mode)")
    ap.add_argument("index", help="index.html 路径")
    ap.add_argument("-o", "--out", default=None, help="输出 .pptx 路径（默认与 index 同名）")
    ap.add_argument("--workers", type=int, default=DEFAULT_WORKERS, help="并行截图进程数")
    args = ap.parse_args()

    index = os.path.abspath(args.index)
    if not os.path.exists(index):
        sys.exit("[html2pptx] 找不到 %s" % index)
    base = os.path.dirname(index)
    out = os.path.abspath(args.out) if args.out else \
        os.path.join(base, re.sub(r"\.html?$", "", os.path.basename(index)) + ".pptx")

    browser = find_browser()
    capture = os.path.join(base, "_html2pptx_capture.html")
    render_dir = os.path.join(base, "_html2pptx_render")

    total = build_capture(index, capture)
    print("[html2pptx] 共 %d 页，浏览器: %s" % (total, os.path.basename(browser)))
    render_all(browser, capture, render_dir, total, args.workers)
    n = build_pptx(render_dir, out, total)
    print("[html2pptx] 完成 -> %s (%d 页)" % (out, n))


if __name__ == "__main__":
    main()