#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
html2pptx_editable — 把 guizang-ppt-skill 的单文件 HTML deck 转成【可编辑】.pptx

原理（无损可编辑的可行路径）:
    headless 浏览器逐页读出每个可见元素的精确几何与样式（坐标/宽高/字号/字重/
    字色/字体/对齐/行高/背景色/圆角/图片），python-pptx 在 16:9 画布的相同坐标
    生成原生文本框、图片、圆角矩形 —— 位置与排版无损，文字全部可编辑。

用法:
    python scripts/html2pptx_editable.py <index.html> [-o 输出.pptx] [--workers 8]

依赖:
    pip install python-pptx
    系统自带 Microsoft Edge 或 Google Chrome（自动探测）

已知边界（非 bug，是方案物理边界）:
    - 字体：PPTX 只写入字体名（如 "Noto Serif SC"），目标机器未安装则由
      PPT 自动 fallback；位置字号颜色不受影响。
    - SVG 图标 / WebGL 流体背景不逐个重建（可编辑版的视觉底图可用图片版补足）。
    - 渐变/阴影/边框细节做近似（纯色块 + 圆角），布局与文字保真。

坑位防御（与图片版同源，勿删）:
    1) 提取前剥离 HTML 注释（模板注释含假 section 标签）
    2) 单页隔离只用 display，禁止 translateX（见 scripts/html2pptx.py 坑2说明）
    3) virtual-time-budget 固定 6000ms
    4) 并行实例独立 --user-data-dir
"""

import argparse
import html as htmllib
import json
import os
import re
import shutil
import subprocess
import sys
import tempfile
import urllib.parse
from concurrent.futures import ThreadPoolExecutor

W, H = 1920, 1080
BUDGET_MS = 6000
DEFAULT_WORKERS = 8
EXTRACT_ID = "__H2P_JSON__"

EXTRACT_TPL = r"""<!DOCTYPE html>
<html lang="zh-CN">
<head>
<meta charset="UTF-8">
<title>h2p-extract</title>
__STYLE__
<style>
  html,body{width:__W__px!important;height:__H__px!important;overflow:hidden}
  .bg{display:none!important}
  #nav,#hint{display:none!important}
  [data-anim]{opacity:1!important;transform:none!important;animation:none!important;transition:none!important}
</style>
</head>
<body>
<div id="deck">__SLIDES__</div>
<pre id="__H2P_JSON__" style="display:none"></pre>
<script>
var q=new URLSearchParams(location.search),
    P=Math.max(1,Math.min(__TOTAL__,parseInt(q.get('p')||'1',10))),
    slides=document.querySelectorAll('#deck .slide');
/* 坑2：只用 display 隔离当前页 */
slides.forEach(function(s,i){s.style.display=(i===P-1)?'':'none'});

(function(){
  var page=slides[P-1], out={page:P,bg:'#ffffff',blocks:[],images:[],texts:[]};

  function px(v){return parseFloat(v)||0}
  function hexNorm(h){
    h=h.replace('#','');
    if(h.length===3)h=h[0]+h[0]+h[1]+h[1]+h[2]+h[2];
    if(h.length===8)h=h.slice(0,6);
    if(h.length!==6||/[^0-9a-fA-F]/.test(h))return null;
    return '#'+h.toLowerCase();
  }
  function rgbToHex(r,g,b){
    return '#'+[r,g,b].map(function(x){return ('0'+Math.round(x).toString(16)).slice(-2)}).join('');
  }
  /* 返回 {hex:'#rrggbb', a:0~1}；仅全透明(≤0.02)视为无色。
     半透明是杂志风 deck 的核心视觉（rgba(ink,.05) 纸感灰底），必须保留 alpha。 */
  function normColor(c){
    if(!c)return null;
    if(c[0]==='#'){var h=hexNorm(c);return h?{hex:h,a:1}:null}
    var m=/rgba?\(([^)]+)\)/.exec(c);
    if(!m)return null;
    var a=m[1].split(',').map(function(x){return parseFloat(x)});
    if(a.length>3&&a[3]<=0.02)return null;
    return {hex:rgbToHex(a[0],a[1],a[2]),a:a.length>3?a[3]:1};
  }
  /* 渐变背景取首色近似（deck 卡片大量使用 linear-gradient） */
  function bgApprox(cs){
    var c=normColor(cs.backgroundColor);
    if(c)return c;
    var bi=cs.backgroundImage||'';
    if(bi&&bi!=='none'){
      var m=/#[0-9a-fA-F]{3,8}|rgba?\([^)]+\)/.exec(bi);
      if(m){var n=normColor(m[0]);if(n)return n}
    }
    return null;
  }
  function radiusOf(el,cs){
    var r=Math.min(px(cs.borderTopLeftRadius),px(cs.borderTopRightRadius),
                   px(cs.borderBottomLeftRadius),px(cs.borderBottomRightRadius));
    return r>2?Math.min(r,60):0;
  }
  function visible(el,cs,rect){
    if(el.closest('#nav,#hint,.bg'))return false;
    return cs.display!=='none'&&cs.visibility!=='hidden'&&rect.width>1&&rect.height>1
      &&rect.bottom>0&&rect.right>0&&rect.top<=__H__&&rect.left<=__W__;
  }

  /* --- 页底色 --- */
  var pcs=getComputedStyle(page);
  var pb=bgApprox(pcs);
  out.bg=pb?pb.hex:(page.classList.contains('dark')?'#0d1b2a':'#ffffff');

  /* --- 背景块 / 图片 / 文本 三类采集 --- */
  var all=page.querySelectorAll('*');
  Array.prototype.forEach.call(all,function(el){
    var cs=getComputedStyle(el), r=el.getBoundingClientRect();
    if(!visible(el,cs,r))return;
    var role=null;
    if(el.tagName==='IMG'&&el.src&&/^file:|^https?:/.test(el.src))role='img';
    /* 背景块：实底色或渐变首色 且 不是含文字的叶子（避免文字底色重复画块） */
    var bg=bgApprox(cs);
    var hasOwnText=Array.prototype.some.call(el.childNodes,function(n){
      return n.nodeType===3&&n.textContent.trim().length>0});
    if(!role&&bg&&r.width*r.height>1200&&!hasOwnText)role='block';
    if(!role&&hasOwnText)role='text';
    /* 带文字的容器（callout/卡片等）底色同样要画块，文字叠其上 */
    if(role==='text'&&bg&&r.width*r.height>1200){
      out.blocks.push({x:px(r.x),y:px(r.y),w:px(r.width),h:px(r.height),
                       color:bg.hex,alpha:bg.a,radius:radiusOf(el,cs)});
    }

    if(role==='block'){
      out.blocks.push({x:px(r.x),y:px(r.y),w:px(r.width),h:px(r.height),
                       color:bg.hex,alpha:bg.a,radius:radiusOf(el,cs)});
    }else if(role==='img'){
      out.images.push({x:px(r.x),y:px(r.y),w:px(r.width),h:px(r.height),
                       src:el.getAttribute('src')});
    }else if(role==='text'){
      var text=Array.prototype.filter.call(el.childNodes,function(n){
        return n.nodeType===3}).map(function(n){return n.textContent}).join('');
      var tc=normColor(cs.color);
      out.texts.push({x:px(r.x),y:px(r.y),w:px(r.width),h:px(r.height),
        size:px(cs.fontSize),weight:cs.fontWeight,
        color:tc?tc.hex:'#111111',alpha:tc?tc.a:1,
        family:(cs.fontFamily||'').split(',')[0].replace(/["']/g,'').trim()||'Calibri',
        align:(cs.textAlign==='start'||cs.textAlign==='end')?'left':cs.textAlign,
        lh:px(cs.lineHeight)/Math.max(1,px(cs.fontSize))||1.2,
        text:text.trim()});
    }
  });

  /* 去噪：丢弃与页面同色的超大背景块 */
  var bg=out.bg.toUpperCase();
  out.blocks=out.blocks.filter(function(b){
    return !(b.color.toUpperCase()===bg&&b.w>=__W__*0.95)});

  document.getElementById('__H2P_JSON__').textContent=JSON.stringify(out);
})();
</script>
</body>
</html>
"""


def find_browser():
    for c in [os.path.expandvars(r"%ProgramFiles(x86)%\Microsoft\Edge\Application\msedge.exe"),
              os.path.expandvars(r"%ProgramFiles%\Microsoft\Edge\Application\msedge.exe"),
              os.path.expandvars(r"%LocalAppData%\Google\Chrome\Application\chrome.exe"),
              os.path.expandvars(r"%ProgramFiles%\Google\Chrome\Application\chrome.exe"),
              os.path.expandvars(r"%ProgramFiles(x86)%\Google\Chrome\Application\chrome.exe"),
              "/usr/bin/google-chrome", "/usr/bin/chromium-browser", "/usr/bin/chromium",
              "/Applications/Google Chrome.app/Contents/MacOS/Google Chrome"]:
        if c and os.path.exists(c):
            return c
    sys.exit("[h2p-editable] 未找到 Edge/Chrome")


def build_extract(index_path, extract_path):
    src = open(index_path, "r", encoding="utf-8").read()
    src = re.sub(r"<!--.*?-->", "", src, flags=re.S)          # 坑1
    styles = re.findall(r"<style[^>]*>.*?</style>", src, re.S)
    slides = re.findall(r'<section class="slide.*?</section>', src, re.S)
    if not styles or not slides:
        sys.exit("[h2p-editable] index.html 结构异常（缺 style 或 slide）")
    tpl = (EXTRACT_TPL.replace("__STYLE__", "\n".join(styles))
           .replace("__SLIDES__", "\n".join(slides))
           .replace("__TOTAL__", str(len(slides)))
           .replace("__H2P_JSON__", EXTRACT_ID)
           .replace("__W__", str(W)).replace("__H__", str(H)))
    open(extract_path, "w", encoding="utf-8").write(tpl)
    return len(slides)


def dump_page(browser, extract_path, p):
    """headless dump-dom，抠出 JSON。返回 dict 或 None。"""
    url = "file:///" + urllib.parse.quote(os.path.abspath(extract_path).replace("\\", "/"))
    prof = os.path.join(tempfile.gettempdir(), "h2pe_%d_%d" % (os.getpid(), p))
    cmd = [browser, "--headless", "--disable-gpu",
           "--user-data-dir=" + prof,
           "--enable-unsafe-swiftshader",
           "--window-size=%d,%d" % (W, H),
           "--virtual-time-budget=%d" % BUDGET_MS,
           "--dump-dom", "%s?p=%d" % (url, p)]
    for _ in range(3):
        try:
            r = subprocess.run(cmd, capture_output=True, timeout=90)
        except subprocess.TimeoutExpired:
            continue
        m = re.search((r'<pre id="%s"[^>]*>(.*?)</pre>' % EXTRACT_ID).encode(),
                      r.stdout, re.S)
        if m:
            try:
                shutil.rmtree(prof, ignore_errors=True)
                return json.loads(htmllib.unescape(m.group(1).decode("utf-8", "ignore")))
            except (ValueError, UnicodeDecodeError):
                pass
    shutil.rmtree(prof, ignore_errors=True)
    return None


def _hex2rgb(h):
    h = h.lstrip("#")
    if len(h) == 3:
        h = "".join(c * 2 for c in h)
    return tuple(int(h[i:i + 2], 16) for i in (0, 2, 4))


def _set_alpha(color_parent, alpha):
    """向 <a:srgbClr> 注入 <a:alpha>（python-pptx 无公开 API，走 lxml）。"""
    if alpha >= 0.99:
        return
    from pptx.oxml.ns import qn as _qn
    srgb = color_parent.find(_qn("a:srgbClr"))
    if srgb is None:
        return
    for old in srgb.findall(_qn("a:alpha")):
        srgb.remove(old)
    el = srgb.makeelement(_qn("a:alpha"), {"val": str(int(alpha * 100000))})
    srgb.append(el)


def _solid(shape, hexcolor, alpha=1.0):
    from pptx.dml.color import RGBColor
    from pptx.oxml.ns import qn as _qn
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*_hex2rgb(hexcolor))
    if alpha < 0.99:
        _set_alpha(shape.fill._xPr.find(_qn("a:solidFill")), alpha)
    shape.line.fill.background()
    shape.shadow.inherit = False


def build_pptx(pages, base_dir, out_path):
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

    SX, SY = 13.333 / W, 7.5 / H           # px -> inch

    def IN(px_, axis=0):
        return Inches(px_ * (SX if axis == 0 else SY))

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    aligns = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
              "right": PP_ALIGN.RIGHT, "justify": PP_ALIGN.JUSTIFY}

    for d in pages:
        s = prs.slides.add_slide(blank)
        # 1) 页底色
        bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        _solid(bg, d.get("bg") or "#ffffff")
        # 2) 背景块（DOM 顺序叠放）
        for b in d.get("blocks", []):
            sh = s.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE if b.get("radius") else MSO_SHAPE.RECTANGLE,
                IN(b["x"]), IN(b["y"]), IN(b["w"]), IN(b["h"]))
            if b.get("radius"):
                try:
                    sh.adjustments[0] = min(0.5, b["radius"] / max(1, min(b["w"], b["h"])))
                except Exception:
                    pass
            _solid(sh, b["color"], b.get("alpha", 1.0))
        # 3) 图片
        for im in d.get("images", []):
            src = im["src"]
            if src and not re.match(r"^(https?:|file:|data:)", src):
                src = os.path.join(base_dir, src)          # 相对路径解析
            if src.startswith("file:///"):
                src = urllib.parse.unquote(src[8:])
            try:
                s.shapes.add_picture(src, IN(im["x"]), IN(im["y"]),
                                      IN(im["w"]), IN(im["h"]))
            except Exception:
                pass
        # 4) 文本（原生可编辑）
        for t in d.get("texts", []):
            tb = s.shapes.add_textbox(IN(t["x"]), IN(t["y"]), IN(t["w"]), IN(t["h"]))
            tf = tb.text_frame
            tf.word_wrap = True
            tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
            tf.vertical_anchor = MSO_ANCHOR.TOP
            para = tf.paragraphs[0]
            para.alignment = aligns.get(t.get("align"), PP_ALIGN.LEFT)
            try:
                para.line_spacing = max(0.8, min(3.0, t.get("lh") or 1.2))
            except Exception:
                pass
            run = para.add_run()
            run.text = t.get("text", "")
            f = run.font
            f.size = Pt(max(6, t.get("size", 16) * 0.75))
            f.name = t.get("family") or "Calibri"
            f.bold = str(t.get("weight", "400")) >= "600"
            try:
                f.color.rgb = RGBColor(*_hex2rgb(t.get("color") or "#111111"))
                a = t.get("alpha", 1.0)
                if a < 0.99:
                    from pptx.oxml.ns import qn as _qn
                    rPr = run._r.get_or_add_rPr()
                    _set_alpha(rPr.find(_qn("a:solidFill")), a)
            except Exception:
                pass
    prs.save(out_path)
    return len(prs.slides._sldIdLst)


def main():
    ap = argparse.ArgumentParser(description="HTML deck -> editable PPTX")
    ap.add_argument("index")
    ap.add_argument("-o", "--out", default=None)
    ap.add_argument("--workers", type=int, default=DEFAULT_WORKERS)
    args = ap.parse_args()

    index = os.path.abspath(args.index)
    if not os.path.exists(index):
        sys.exit("[h2p-editable] 找不到 %s" % index)
    base = os.path.dirname(index)
    out = os.path.abspath(args.out) if args.out else \
        os.path.join(base, re.sub(r"\.html?$", "", os.path.basename(index)) + "_editable.pptx")

    browser = find_browser()
    extract = os.path.join(base, "_h2p_extract.html")
    total = build_extract(index, extract)
    print("[h2p-editable] 共 %d 页，浏览器: %s" % (total, os.path.basename(browser)))

    with ThreadPoolExecutor(max_workers=args.workers) as ex:
        dumps = list(ex.map(lambda p: dump_page(browser, extract, p),
                            range(1, total + 1)))
    pages = [d for d in dumps if d]
    if len(pages) != total:
        sys.exit("[h2p-editable] 提取失败 %d 页（重试3次仍无效）" % (total - len(pages)))

    n = build_pptx(pages, base, out)
    print("[h2p-editable] 完成 -> %s (%d 页，全文字可编辑)" % (out, n))


if __name__ == "__main__":
    main()